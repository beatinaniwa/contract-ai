"""Utilities to load plain text from uploaded files."""

from __future__ import annotations

import io
from collections.abc import Iterable
from pathlib import Path

from pypdf import PdfReader
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.exc import PackageNotFoundError

_TEXT_EXTENSIONS = {".txt", ".md", ".markdown"}


def load_text_from_bytes(data: bytes, filename: str) -> str:
    """Extract text content from raw file bytes.

    Args:
        data: Raw file content.
        filename: Original file name used for extension detection.

    Returns:
        Extracted text stripped of leading/trailing whitespace.

    Raises:
        ValueError: If the file type is unsupported or text cannot be extracted.
    """

    if not data:
        raise ValueError("ファイルの内容が空です。")

    suffix = Path(filename or "uploaded").suffix.lower()

    if suffix in _TEXT_EXTENSIONS:
        return _decode_text_file(data)

    if suffix == ".pdf":
        return _extract_pdf_text(data)

    if suffix == ".pptx":
        return _extract_pptx_text(data)

    raise ValueError(f"サポートされていないファイル形式です: {suffix or '不明'}")


def _decode_text_file(data: bytes) -> str:
    text = data.decode("utf-8", errors="ignore").strip()
    if not text:
        raise ValueError("テキストファイルから内容を取得できませんでした。")
    return text


def _extract_pdf_text(data: bytes) -> str:
    reader = PdfReader(io.BytesIO(data))
    extracted_parts = []
    for page in reader.pages:
        page_text = page.extract_text() or ""
        if page_text:
            extracted_parts.append(page_text.strip())
    combined = "\n\n".join(part for part in extracted_parts if part)
    combined = combined.strip()
    if not combined:
        raise ValueError("PDFからテキストを抽出できませんでした。")
    return combined


def _extract_pptx_text(data: bytes) -> str:
    try:
        presentation = Presentation(io.BytesIO(data))
    except PackageNotFoundError as exc:
        raise ValueError("PPTXファイルが破損している可能性があります。") from exc
    except Exception as exc:  # pragma: no cover - python-pptx internal errors are rare
        raise ValueError("PPTXファイルの読み込みに失敗しました。") from exc

    slides_output: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        lines = list(_iter_slide_text(slide))
        if not lines:
            continue
        bullets = "\n".join(f"- {line}" for line in _deduplicate_preserving_order(lines))
        slides_output.append(f"[Slide {index}]\n{bullets}")

    if not slides_output:
        raise ValueError("PPTXからテキストを抽出できませんでした。")

    return "\n\n".join(slides_output)


def _iter_slide_text(slide) -> Iterable[str]:
    for shape in slide.shapes:
        yield from _iter_shape_text(shape)


def _iter_shape_text(shape) -> Iterable[str]:
    shape_type = getattr(shape, "shape_type", None)
    if shape_type == MSO_SHAPE_TYPE.GROUP:
        for subshape in shape.shapes:
            yield from _iter_shape_text(subshape)
        return

    if getattr(shape, "has_text_frame", False):
        text_frame = shape.text_frame
        for paragraph in getattr(text_frame, "paragraphs", []):
            text = "".join(run.text for run in getattr(paragraph, "runs", [])).strip()
            if text:
                for line in text.splitlines():
                    cleaned = line.strip()
                    if cleaned:
                        yield cleaned
        return

    if getattr(shape, "has_table", False):
        table = shape.table
        for row in table.rows:
            for cell in row.cells:
                for line in cell.text.splitlines():
                    cleaned = line.strip()
                    if cleaned:
                        yield cleaned
        return


def _deduplicate_preserving_order(lines: list[str]) -> list[str]:
    seen: set[str] = set()
    deduped: list[str] = []
    for line in lines:
        normalized = " ".join(line.split())
        if normalized and normalized not in seen:
            deduped.append(normalized)
            seen.add(normalized)
    return deduped

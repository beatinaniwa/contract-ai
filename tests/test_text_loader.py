import io

import pytest
from pptx import Presentation

from services.text_loader import load_text_from_bytes


def _make_pdf_with_text(text: str) -> bytes:
    def _escape(content: str) -> str:
        return content.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")

    escaped = _escape(text)
    content = f"BT\n/F1 24 Tf\n72 712 Td\n({escaped}) Tj\nET\n"
    content_bytes = content.encode("utf-8")
    objects = [
        "<< /Type /Catalog /Pages 2 0 R >>",
        "<< /Type /Pages /Count 1 /Kids [3 0 R] >>",
        "<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>",
        f"<< /Length {len(content_bytes)} >>\nstream\n{content}\nendstream",
        "<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]

    buffer = io.BytesIO()
    buffer.write(b"%PDF-1.4\n")
    offsets = [0]
    for index, obj in enumerate(objects, start=1):
        offsets.append(buffer.tell())
        buffer.write(f"{index} 0 obj\n".encode("ascii"))
        buffer.write(obj.encode("utf-8"))
        buffer.write(b"\nendobj\n")

    xref_pos = buffer.tell()
    buffer.write(f"xref\n0 {len(offsets)}\n".encode("ascii"))
    buffer.write(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        buffer.write(f"{offset:010d} 00000 n \n".encode("ascii"))
    buffer.write(b"trailer\n")
    buffer.write(f"<< /Size {len(offsets)} /Root 1 0 R >>\n".encode("ascii"))
    buffer.write(b"startxref\n")
    buffer.write(f"{xref_pos}\n".encode("ascii"))
    buffer.write(b"%%EOF")
    return buffer.getvalue()


def _make_pptx_with_slides(slides: list[tuple[str, list[str]]]) -> bytes:
    presentation = Presentation()
    layout = presentation.slide_layouts[1]  # Title and Content
    for title, bullet_points in slides:
        slide = presentation.slides.add_slide(layout)
        slide.shapes.title.text = title
        body = slide.placeholders[1]
        body.text = "\n".join(bullet_points)
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def test_load_text_from_txt_bytes():
    content = "案件概要を記載したサンプルテキストです。\n詳細条件も含まれます。"
    result = load_text_from_bytes(content.encode("utf-8"), "sample.txt")
    assert "案件概要" in result
    assert "詳細条件" in result


def test_load_text_from_pdf_bytes():
    pdf_bytes = _make_pdf_with_text("Hello PDF")
    result = load_text_from_bytes(pdf_bytes, "sample.pdf")
    assert result == "Hello PDF"


def test_load_text_from_bytes_raises_for_unknown_extension():
    with pytest.raises(ValueError):
        load_text_from_bytes(b"binary", "sample.docx")


def test_load_text_from_bytes_raises_for_empty_payload():
    with pytest.raises(ValueError):
        load_text_from_bytes(b"", "empty.txt")


def test_load_text_from_pptx_bytes_extracts_slide_text():
    pptx_bytes = _make_pptx_with_slides(
        [
            ("提案概要", ["条件A", "条件B"]),
            ("スケジュール", ["開始日: 4/1"]),
        ]
    )

    result = load_text_from_bytes(pptx_bytes, "slides.pptx")

    assert "[Slide 1]" in result
    assert "- 提案概要" in result
    assert "- 条件A" in result
    assert "[Slide 2]" in result
    assert "- スケジュール" in result


def test_load_text_from_pptx_bytes_raises_for_invalid_file():
    with pytest.raises(ValueError):
        load_text_from_bytes(b"not-a-pptx", "slides.pptx")
